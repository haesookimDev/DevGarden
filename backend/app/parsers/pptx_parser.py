import io

from pptx import Presentation

from app.parsers.base import DocumentParser, ParsedDocument


class PPTXParser(DocumentParser):
    def parse(self, file_data: bytes) -> ParsedDocument:
        prs = Presentation(io.BytesIO(file_data))
        text_parts = []
        images = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- Slide {slide_num} ---"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_texts.append(paragraph.text)

                if shape.shape_type == 13:  # Picture
                    images.append(shape.image.blob)

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker Notes]: {notes}")

            text_parts.append("\n".join(slide_texts))

        return ParsedDocument(
            text="\n\n".join(text_parts),
            images=images,
            metadata={},
            page_count=len(prs.slides),
        )
